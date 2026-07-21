"""Best-effort text extraction from common file types."""

from __future__ import annotations

import json
import logging
import math
import re
import shutil
import subprocess
import sys
import threading
import unicodedata
from html.parser import HTMLParser
from io import BytesIO
from pathlib import Path

_log = logging.getLogger(__name__)


_FITZ_BUNDLE_SCRIPT = """\
import sys, json, fitz
path = sys.argv[1]
try:
    doc = fitz.open(path)
    n = len(doc)
    flags = (int(fitz.TEXT_PRESERVE_WHITESPACE) | int(fitz.TEXT_MEDIABOX_CLIP)
             | int(fitz.TEXT_DEHYPHENATE) | int(fitz.TEXT_USE_CID_FOR_UNKNOWN_UNICODE))
    pages = [doc[i].get_text('text', flags=flags) or '' for i in range(n)]
    xrefs = {img[0] for i in range(n) for img in doc[i].get_images(full=True)}
    has_imgs = any(doc[i].get_images(full=True) for i in range(min(3, n)))
    doc.close()
    print(json.dumps({'ok': True, 'pages': pages, 'page_count': n,
                      'image_xrefs': len(xrefs), 'has_images': bool(has_imgs)}))
except Exception as e:
    print(json.dumps({'ok': False, 'error': str(e)}))
"""

_FITZ_EMPTY = {"ok": False, "pages": [], "page_count": None, "image_xrefs": 0, "has_images": False}

# Thread-local cache: avoid re-spawning a subprocess for the same file within one thread.
# All fitz calls for a given file happen sequentially in the same worker thread,
# so caching the last result eliminates 2 of 3 subprocess spawns per PDF.
_fitz_tls = threading.local()


def _fitz_bundle(path: Path, timeout: int = 120) -> dict:
    """Run all fitz operations for ``path`` in one subprocess and cache per thread.

    Returns a dict with keys: ok, pages, page_count, image_xrefs, has_images.
    """
    path_str = str(path)
    if getattr(_fitz_tls, "path", None) == path_str:
        return _fitz_tls.result  # type: ignore[return-value]
    try:
        proc = subprocess.run(
            [sys.executable, "-c", _FITZ_BUNDLE_SCRIPT, path_str],
            capture_output=True, text=True, timeout=timeout,
        )
        if proc.returncode == 0 and proc.stdout.strip():
            result = json.loads(proc.stdout)
        else:
            result = dict(_FITZ_EMPTY)
    except Exception:
        result = dict(_FITZ_EMPTY)
    _fitz_tls.path = path_str
    _fitz_tls.result = result
    return result


MAX_CHARS_DEFAULT = 4000
MAX_BYTES_READ_DEFAULT = 2_000_000  # 2 MB safety cap for text-like reads
# Full-document extraction for word counts (Notion-sized exports; very large files may be slow).
MAX_BYTES_READ_FULL = 100 * 1024 * 1024  # 100 MB

_WORD_TOKENS = re.compile(r"\S+", re.UNICODE)

# Video containers: we do not read file bytes for classification — filename stem is the "title".
VIDEO_FILE_EXTENSIONS = frozenset(
    {"mp4", "mov", "m4v", "avi", "mkv", "webm", "wmv", "flv", "mpeg", "mpg", "3gp"}
)

# Audio containers: same policy as video — no decode/transcribe; classify from name/path only.
AUDIO_FILE_EXTENSIONS = frozenset(
    {"mp3", "wav", "m4a", "aac", "flac", "ogg", "opus", "wma", "aiff", "aif", "alac", "caf"}
)


def is_video_file(path: str | Path) -> bool:
    """True if ``path`` is a video container by suffix (no content sniffing)."""
    ext = Path(path).suffix.lower().lstrip(".")
    return ext in VIDEO_FILE_EXTENSIONS


def is_audio_file(path: str | Path) -> bool:
    """True if ``path`` is an audio container by suffix (no content sniffing)."""
    ext = Path(path).suffix.lower().lstrip(".")
    return ext in AUDIO_FILE_EXTENSIONS


def is_title_only_media(path: str | Path) -> bool:
    """Video or audio: do not read bytes for classification — use filename stem + path only."""
    return is_video_file(path) or is_audio_file(path)


def media_title_text(path: str | Path) -> str:
    """Filename stem used as the only text signal for video/audio artifacts."""
    p = Path(path)
    s = p.stem.strip()
    return s if s else p.name


def video_title_text(path: str | Path) -> str:
    """Alias of :func:`media_title_text` (backward compatible)."""
    return media_title_text(path)


_IMAGE_DOT_EXT = frozenset({".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"})
_MAX_IMAGE_OCR_PIXELS = 12_000_000


def _resize_for_ocr(im: object) -> object:
    from PIL import Image as PILImage

    if not isinstance(im, PILImage.Image):
        return im
    w, h = im.size
    pixels = w * h
    if pixels <= _MAX_IMAGE_OCR_PIXELS:
        return im
    scale = math.sqrt(_MAX_IMAGE_OCR_PIXELS / float(max(1, pixels)))
    nw = max(1, int(w * scale))
    nh = max(1, int(h * scale))
    try:
        resample = PILImage.Resampling.LANCZOS  # type: ignore[attr-defined]
    except Exception:
        resample = PILImage.LANCZOS  # type: ignore[attr-defined]
    return im.resize((nw, nh), resample)


def _exif_descriptive_strings(im: object) -> list[str]:
    """Pull a few human-edited EXIF fields when present."""
    parts: list[str] = []
    try:
        from PIL.ExifTags import TAGS

        exif = im.getexif()
        if not exif:
            return parts
        for k, v in exif.items():
            tag = TAGS.get(k, str(k))
            if tag not in (
                "UserComment",
                "ImageDescription",
                "DocumentName",
                "Artist",
                "Copyright",
            ):
                continue
            if isinstance(v, bytes):
                try:
                    v = v.decode("utf-8", errors="replace")
                except Exception:
                    v = str(v)
            if isinstance(v, str) and v.strip():
                parts.append(f"{tag}: {v.strip()}")
    except Exception:
        pass
    return parts


def _extract_image_plaintext(path: Path) -> str | None:
    """Raster / common image formats → plain text: dimensions + EXIF notes + optional OCR."""
    try:
        from PIL import Image
    except Exception:
        return None
    ext = path.suffix.lower()
    if ext not in _IMAGE_DOT_EXT:
        return None
    im = None
    try:
        im = Image.open(path)
        im.load()
        w, h = im.size
        lines: list[str] = [f"[image] {path.name} dimensions={w}x{h} mode={im.mode}"]
        lines.extend(_exif_descriptive_strings(im))

        ocr_piece = ""
        try:
            import pytesseract

            if shutil.which("tesseract"):
                if im.mode in ("RGB", "L"):
                    rgb = im
                else:
                    rgb = im.convert("RGB")
                work = _resize_for_ocr(rgb)
                ocr_piece = (pytesseract.image_to_string(work, lang="eng") or "").strip()
        except Exception:
            ocr_piece = ""

        if ocr_piece:
            lines.append("--- OCR ---")
            lines.append(ocr_piece)

        txt = "\n".join(lines).strip()
        return _normalize_extracted_text(txt) if txt else None
    except Exception:
        return None
    finally:
        if im is not None:
            try:
                im.close()
            except Exception:
                pass


def _rfc822_header_bytes(path: Path, *, max_scan: int = 524_288) -> bytes:
    """Read up to ``max_scan`` bytes and trim to the end of RFC822/MIME headers."""
    try:
        data = path.read_bytes()[:max_scan]
    except OSError:
        return b""
    for sep in (b"\r\n\r\n", b"\n\n"):
        j = data.find(sep)
        if j != -1:
            return data[:j]
    return data


def extract_email_subject(path: str | Path) -> str:
    """Decode ``Subject`` from RFC822/MIME headers only (no message body read)."""
    path = Path(path)
    if path.suffix.lower() != ".eml":
        return ""
    try:
        import email.policy
        from email import message_from_bytes
        from email.header import decode_header, make_header

        hdr = _rfc822_header_bytes(path)
        if not hdr.strip():
            return ""
        msg = message_from_bytes(hdr + b"\n", policy=email.policy.default)
        raw = msg.get("Subject")
        if raw is None:
            return ""
        s = str(raw).strip()
        if not s:
            return ""
        try:
            return str(make_header(decode_header(s))).strip()
        except Exception:
            return s
    except Exception:
        return ""


# GPT-family-style tokenization (also used as a generic LLM planning estimate).
_LLM_TOKEN_ENCODING_NAME = "cl100k_base"
_tiktoken_enc: object = None
_tiktoken_lock = threading.Lock()


def _get_tiktoken_enc() -> object:
    global _tiktoken_enc
    if _tiktoken_enc is not None:
        return _tiktoken_enc
    with _tiktoken_lock:
        if _tiktoken_enc is None:
            import tiktoken
            _tiktoken_enc = tiktoken.get_encoding(_LLM_TOKEN_ENCODING_NAME)
    return _tiktoken_enc

# Heuristic divisor when a format has no native page count.
ESTIMATED_WORDS_PER_PAGE = 275
# Notion/HTML→docx-style documents: ``docProps/app.xml`` ``Pages`` from converters is
# unreliable; use a lower words/page estimate so totals align better with Word layout.
ESTIMATED_WORDS_PER_PAGE_DOCUMENT = 120
# DOCX with **no** inline shapes (pure prose): words/page varies—short pages read
# denser than long Notion→Word exports (looser spacing). Use tiered divisors below.
ESTIMATED_WORDS_PER_PAGE_DOCUMENT_TEXT_ONLY = 141
ESTIMATED_WORDS_PER_PAGE_DOCUMENT_TEXT_ONLY_LONG = 180
DOCX_TEXT_ONLY_LONG_WORD_THRESHOLD = 4000
# Sum of ``inline_shapes`` heights (EMU) / this value → "stack pages" from diagrams /
# screenshots beyond text-density. Calibrated so typical Notion→docx exports match
# Word more closely when figures consume vertical space words alone undercount.
DOCX_IMAGE_STACK_PAGE_EMU = 12_374_880


_EMOJI_RE = re.compile(
    "["
    "\U0001F1E6-\U0001F1FF"  # flags
    "\U0001F300-\U0001FAFF"  # emoji & symbols (broad)
    "\U00002700-\U000027BF"  # dingbats
    "\U00002600-\U000026FF"  # misc symbols
    "]+",
    flags=re.UNICODE,
)


def _truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "\n…(truncated)…"


def _normalize_extracted_text(text: str) -> str:
    """Normalize common extraction artifacts (PDFs/OCR/etc.)."""
    if not text:
        return text
    # Strip emoji/pictographs; they often appear in doc titles and look noisy in summaries.
    text = _EMOJI_RE.sub(" ", text)
    # Expand compatibility ligatures (ﬁ -> fi, ﬂ -> fl, etc.) for cleaner words.
    text = unicodedata.normalize("NFKC", text)
    # Drop emoji variation selectors that can leave stray glyphs.
    text = text.replace("\ufe0f", "")
    # Fix hyphenation across line breaks: "exam-\nple" -> "example"
    text = re.sub(r"(\w)-\n(\w)", r"\1\2", text)
    # Fix word-wrapping splits without hyphenation (common in extracted PDFs):
    # "inven\ntory" -> "inventory"
    text = re.sub(r"([a-z])\n([a-z])", r"\1\2", text)
    text = re.sub(r"([A-Z])\n([A-Z])", r"\1\2", text)
    # Convert newlines to spaces (later callers may re-split if needed)
    text = text.replace("\r", "\n")
    # Remove soft hyphen and zero-width spaces
    text = text.replace("\u00ad", "").replace("\u200b", "")
    # Normalize common ligatures
    text = (
        text.replace("\ufb01", "fi")
        .replace("\ufb02", "fl")
        .replace("\ufb03", "ffi")
        .replace("\ufb04", "ffl")
    )
    return text


def _pdf_text_quality_score(text: str) -> float:
    """Prefer extracts with more letters and fewer single-letter word fragments."""
    if not text or not text.strip():
        return -1.0
    letters = sum(1 for c in text if c.isalpha())
    toks = text.split()
    if not toks:
        return letters / max(len(text), 1)
    singles = sum(1 for t in toks if len(t) == 1 and t.isalpha())
    frag = singles / len(toks)
    return (letters / max(len(text), 1)) * 1.5 - frag


def _extract_pdf_pages(path: Path, *, max_pages: int | None) -> str | None:
    """Extract PDF text using PyMuPDF first; pypdf only as fallback for blank pages."""
    bundle = _fitz_bundle(path)
    fitz_pages: list[str] = []
    fitz_ok = bundle["ok"]
    blank_indices: list[int] = []

    if fitz_ok:
        raw_pages: list[str] = bundle["pages"]
        if max_pages is not None:
            raw_pages = raw_pages[:max_pages]
        for i, t in enumerate(raw_pages):
            cleaned = _normalize_extracted_text(t) if t.strip() else ""
            fitz_pages.append(cleaned)
            if not cleaned.strip():
                blank_indices.append(i)

    # Only run pypdf for pages fitz left blank (or for the whole doc if fitz failed).
    pypdf_by_index: dict[int, str] = {}
    need_pypdf = set(blank_indices) if fitz_ok else None  # None = all pages
    if not fitz_ok or blank_indices:
        try:
            from pypdf import PdfReader  # type: ignore

            reader = PdfReader(str(path))
            n_all = len(reader.pages)
            n = n_all if max_pages is None else min(max_pages, n_all)
            for i in range(n):
                if need_pypdf is None or i in need_pypdf:
                    t = reader.pages[i].extract_text() or ""
                    cleaned = _normalize_extracted_text(t) if t.strip() else ""
                    if cleaned.strip():
                        pypdf_by_index[i] = cleaned
        except Exception:
            _log.debug("pypdf extraction failed for %s", path, exc_info=True)

    if not fitz_pages and not pypdf_by_index:
        return None

    n = len(fitz_pages) if fitz_pages else (max(pypdf_by_index.keys()) + 1 if pypdf_by_index else 0)
    merged: list[str] = []
    for i in range(n):
        page_text = fitz_pages[i] if i < len(fitz_pages) else ""
        if not page_text.strip() and i in pypdf_by_index:
            page_text = pypdf_by_index[i]
        if page_text.strip():
            merged.append(page_text)

    return "\n\n".join(merged) if merged else None


def _read_text_file(path: Path, *, max_bytes: int) -> str:
    data = path.read_bytes()[:max_bytes]
    return data.decode("utf-8", errors="replace")


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._skip_depth = 0  # inside script/style
        self._chunks: list[str] = []

    def handle_starttag(self, tag: str, attrs) -> None:  # type: ignore[override]
        t = (tag or "").lower()
        if t in {"script", "style"}:
            self._skip_depth += 1

    def handle_endtag(self, tag: str) -> None:  # type: ignore[override]
        t = (tag or "").lower()
        if t in {"script", "style"} and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data: str) -> None:  # type: ignore[override]
        if self._skip_depth > 0:
            return
        if data and data.strip():
            self._chunks.append(data)

    def text(self) -> str:
        joined = " ".join(self._chunks)
        joined = re.sub(r"[ \t\r\f\v]+", " ", joined)
        return joined.strip()


def _strip_html(html: str) -> str:
    # Regex-stripping HTML is fragile; use a real parser to avoid dropping characters.
    try:
        parser = _HTMLTextExtractor()
        parser.feed(html)
        parser.close()
        return parser.text()
    except Exception:
        # Fallback: best-effort regex strip.
        html = re.sub(r"(?is)<(script|style).*?>.*?</\\1>", " ", html)
        html = re.sub(r"(?is)<[^>]+>", " ", html)
        html = re.sub(r"[ \t\r\f\v]+", " ", html)
        return html.strip()


def _extract_docx_document(document) -> str | None:
    """Plain text from an open ``python-docx`` document (paragraphs + tables)."""
    parts: list[str] = []

    def walk_table(table) -> None:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    if p.text.strip():
                        parts.append(_normalize_extracted_text(p.text))
                for nested in cell.tables:
                    walk_table(nested)

    for p in document.paragraphs:
        if p.text.strip():
            parts.append(_normalize_extracted_text(p.text))
    for table in document.tables:
        walk_table(table)
    return "\n".join(parts) if parts else None


def _extract_docx_full_text_from_bytes(data: bytes) -> str | None:
    try:
        import docx  # type: ignore
    except Exception:
        return None
    try:
        document = docx.Document(BytesIO(data))
    except Exception:
        return None
    return _extract_docx_document(document)


def _extract_docx_full_text(path: Path) -> str | None:
    try:
        import docx  # type: ignore
    except Exception:
        return None
    try:
        document = docx.Document(str(path))
    except Exception:
        return None
    return _extract_docx_document(document)


def _extract_full_text_html_via_docx(path: Path, *, max_bytes: int) -> str | None:
    """HTML → ``html2docx`` → extract docx text (same pipeline as ``tools/html_to_docx``)."""
    try:
        from extractors.html_to_docx_convert import html_string_to_docx_bytes
    except ImportError:
        return None
    try:
        raw_html = _read_text_file(path, max_bytes=max_bytes)
        buf = html_string_to_docx_bytes(
            raw_html,
            title=None,
            raw=False,
            fallback_title=path.stem,
            base_dir=path.parent,
        )
        return _extract_docx_full_text_from_bytes(buf.getvalue())
    except Exception:
        return None


def _extract_xlsx_full_text(path: Path) -> str | None:
    try:
        import openpyxl  # type: ignore
    except Exception:
        return None
    parts: list[str] = []
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    try:
        for sheetname in wb.sheetnames:
            ws = wb[sheetname]
            sheet_rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [("" if v is None else str(v)).strip() for v in row]
                if any(cells):
                    sheet_rows.append("\t".join(cells))
            if sheet_rows:
                parts.append(f"[Sheet: {sheetname}]")
                parts.extend(sheet_rows)
                parts.append("")
    finally:
        wb.close()
    txt = "\n".join(parts).strip()
    return _normalize_extracted_text(txt) if txt else None


def _extract_xls_full_text(path: Path) -> str | None:
    try:
        import xlrd  # type: ignore
    except Exception:
        return None
    try:
        wb = xlrd.open_workbook(str(path))
    except Exception:
        return None
    parts: list[str] = []
    for sheet_idx in range(wb.nsheets):
        ws = wb.sheet_by_index(sheet_idx)
        sheet_rows: list[str] = []
        for row_idx in range(ws.nrows):
            cells = []
            for col_idx in range(ws.ncols):
                ctype = ws.cell_type(row_idx, col_idx)
                val = ws.cell_value(row_idx, col_idx)
                if ctype == xlrd.XL_CELL_EMPTY or ctype == xlrd.XL_CELL_ERROR:
                    cells.append("")
                elif ctype == xlrd.XL_CELL_DATE:
                    try:
                        import datetime
                        dt = xlrd.xldate_as_datetime(val, wb.datemode)
                        cells.append(dt.strftime("%Y-%m-%d"))
                    except Exception:
                        cells.append(str(val))
                elif ctype == xlrd.XL_CELL_NUMBER:
                    cells.append(str(int(val)) if isinstance(val, float) and val == int(val) else str(val))
                elif ctype == xlrd.XL_CELL_BOOLEAN:
                    cells.append("True" if val else "False")
                else:
                    cells.append(str(val).strip())
            if any(c.strip() for c in cells):
                sheet_rows.append("\t".join(cells))
        if sheet_rows:
            parts.append(f"[Sheet: {ws.name}]")
            parts.extend(sheet_rows)
            parts.append("")
    txt = "\n".join(parts).strip()
    return _normalize_extracted_text(txt) if txt else None


def _extract_pst_text(path: Path, *, max_chars: int = MAX_CHARS_DEFAULT) -> str | None:
    """Extract a representative text sample from an Outlook PST/OST archive.

    Walks all mail folders and emits one header+body block per message until
    ``max_chars`` is reached.  Requires ``libpff-python`` (pip install libpff-python).
    """
    try:
        import pypff  # type: ignore
    except ImportError:
        print(f"[pst] libpff-python not installed — cannot read {path.name}", flush=True)
        return None

    def _decode(raw: bytes | str | None) -> str:
        if raw is None:
            return ""
        if isinstance(raw, bytes):
            for enc in ("utf-8", "latin-1", "cp1252"):
                try:
                    return raw.decode(enc)
                except Exception:
                    pass
            return raw.decode("utf-8", errors="replace")
        return str(raw)

    def _rtf_to_text(raw: bytes | str | None) -> str:
        """Strip RTF markup to plain text as a last-resort body fallback."""
        text = _decode(raw)
        if not text:
            return ""
        import re as _re
        text = _re.sub(r"\\[a-z]+[-\d]*\s?", " ", text)
        text = _re.sub(r"[{}\\]", " ", text)
        return _re.sub(r"\s+", " ", text).strip()

    parts: list[str] = []
    total = 0

    def _walk(folder) -> None:
        nonlocal total
        if total >= max_chars:
            return
        try:
            n_msgs = folder.number_of_sub_messages
        except Exception:
            n_msgs = 0
        for i in range(n_msgs):
            if total >= max_chars:
                break
            try:
                msg = folder.get_sub_message(i)
                subject = _decode(msg.subject).strip()
                sender  = _decode(msg.sender_name).strip()
                date    = str(msg.delivery_time or msg.client_submit_time or "")
                body    = _decode(msg.plain_text_body).strip()
                if not body:
                    body = _strip_html(_decode(msg.html_body)).strip()
                if not body:
                    body = _rtf_to_text(msg.rtf_body).strip()
                n_att = msg.number_of_attachments
                att_note = f"  [{n_att} attachment(s)]" if n_att else ""
                block = (
                    f"Subject: {subject}\n"
                    f"From: {sender}  Date: {date}{att_note}\n"
                    f"{body[:600]}\n---"
                )
                parts.append(block)
                total += len(block)
            except Exception as exc:
                _log.debug("PST message %d failed in %s: %s", i, path.name, exc)

        try:
            n_folders = folder.number_of_sub_folders
        except Exception:
            n_folders = 0
        for i in range(n_folders):
            if total >= max_chars:
                break
            try:
                _walk(folder.get_sub_folder(i))
            except Exception as exc:
                _log.debug("PST sub-folder %d failed in %s: %s", i, path.name, exc)

    try:
        pst = pypff.file()
        pst.open(str(path))
        root = pst.get_root_folder()
        if root is None:
            print(f"[pst] WARNING: no root folder found in {path.name}", flush=True)
            pst.close()
            return None
        _walk(root)
        pst.close()
    except Exception as exc:
        print(f"[pst] WARNING: failed to open/read {path.name} — {exc}", flush=True)
        return None

    text = "\n\n".join(parts)
    if not text.strip():
        print(f"[pst] WARNING: no message text extracted from {path.name}", flush=True)
        return None
    return _normalize_extracted_text(text[:max_chars])


def extract_full_text(
    file_path: str | Path,
    *,
    max_bytes: int = MAX_BYTES_READ_FULL,
    html_via_docx: bool = True,
) -> str | None:
    """Return normalized plain text for the whole file (no snippet truncation).

    Subject to ``max_bytes`` for text-like formats; PDFs use all pages. Binary /
    unknown extensions return ``None``.

    For ``.eml`` (RFC822 email), only the **Subject** header is parsed from the first
    part of the file (headers section); the message body is not read.

    For ``.html`` / ``.htm``, when ``html_via_docx`` is True (default), text is taken
    from an in-memory Word conversion (requires ``html2docx``); otherwise tags are
    stripped with :func:`_strip_html`.
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    try:
        if ext == ".eml":
            subj = extract_email_subject(path)
            if not subj:
                return None
            t = _normalize_extracted_text(subj)
            return t if t and t.strip() else None

        if ext in {".txt", ".md", ".rst", ".log"}:
            raw = _read_text_file(path, max_bytes=max_bytes)
            return _normalize_extracted_text(raw) if raw.strip() else None

        if ext in {".html", ".htm"}:
            if html_via_docx:
                t = _extract_full_text_html_via_docx(path, max_bytes=max_bytes)
                if t is not None and t.strip():
                    return t
            raw = _read_text_file(path, max_bytes=max_bytes)
            cleaned = _strip_html(raw)
            return _normalize_extracted_text(cleaned) if cleaned.strip() else None

        if ext == ".json":
            raw = _read_text_file(path, max_bytes=max_bytes)
            try:
                obj = json.loads(raw)
                pretty = json.dumps(obj, ensure_ascii=False, indent=2)
                return _normalize_extracted_text(pretty)
            except Exception:
                return _normalize_extracted_text(raw) if raw.strip() else None

        if ext in {".csv", ".tsv"}:
            raw = _read_text_file(path, max_bytes=max_bytes)
            return _normalize_extracted_text(raw) if raw.strip() else None

        if ext in {".xlsx", ".xlsm"}:
            return _extract_xlsx_full_text(path)

        if ext == ".xls":
            return _extract_xls_full_text(path)

        if ext == ".pdf":
            raw = _extract_pdf_pages(path, max_pages=None)
            return _normalize_extracted_text(raw) if raw and raw.strip() else None

        if ext == ".docx":
            t = _extract_docx_full_text(path)
            return t if t else None

        if ext == ".pptx":
            try:
                from pptx import Presentation  # type: ignore
            except Exception:
                return None
            prs = Presentation(str(path))
            slide_parts: list[str] = []
            for i in range(len(prs.slides)):
                slide = prs.slides[i]
                for shape in slide.shapes:
                    text = getattr(shape, "text", None)
                    if isinstance(text, str) and text.strip():
                        slide_parts.append(_normalize_extracted_text(text.strip()))
            joined = "\n\n".join(slide_parts)
            return _normalize_extracted_text(joined) if joined.strip() else None

        if ext in _IMAGE_DOT_EXT:
            t = _extract_image_plaintext(path)
            if not t:
                return None
            if len(t) > max_bytes:
                t = t[:max_bytes]
            return _normalize_extracted_text(t) if t.strip() else None

        if ext in {".pst", ".ost"}:
            return _extract_pst_text(path, max_chars=max_bytes)

        return None
    except Exception:
        _log.debug("extract_full_text failed for %s", file_path, exc_info=True)
        return None


def count_words(text: str) -> int:
    """Count tokens as non-whitespace runs after ``_normalize_extracted_text``."""
    if not text or not str(text).strip():
        return 0
    normalized = _normalize_extracted_text(str(text))
    if not normalized.strip():
        return 0
    return len(_WORD_TOKENS.findall(normalized))


def count_words_from_file(
    file_path: str | Path,
    *,
    html_via_docx: bool = True,
) -> int:
    """Full-file word count; returns ``0`` if unreadable or empty.

    HTML uses the same conversion path as ``tools/html_to_docx.py`` when
    ``html_via_docx`` is True (default).
    """
    try:
        t = extract_full_text(file_path, html_via_docx=html_via_docx)
        if not t:
            return 0
        return count_words(t)
    except Exception:
        return 0


def count_llm_tokens(text: str) -> int:
    """Subword token count for the **exact extracted document text** (LLM input-sized).

    Uses ``tiktoken`` ``cl100k_base`` (GPT-3.5/4–style). If ``tiktoken`` is unavailable,
    approximates with ``len(text) // 4``. Provider-specific tokenizers (e.g. Anthropic)
    differ slightly—use for capacity planning, not billing reconciliation.
    """
    if not text or not str(text).strip():
        return 0
    raw = str(text)
    try:
        enc = _get_tiktoken_enc()
        return len(enc.encode(raw))
    except Exception:
        return max(0, len(raw) // 4)


def token_count_from_file(
    file_path: str | Path,
    *,
    html_via_docx: bool = True,
) -> int:
    """Same extraction pipeline as :func:`count_words_from_file`; returns LLM-style tokens."""
    try:
        t = extract_full_text(file_path, html_via_docx=html_via_docx)
        if not t:
            return 0
        return count_llm_tokens(t)
    except Exception:
        return 0


def _pages_from_word_estimate(word_count: int, *, per_page: int | None = None) -> int:
    if word_count <= 0:
        return 0
    n = per_page if per_page is not None else ESTIMATED_WORDS_PER_PAGE
    return max(1, (word_count + n - 1) // n)


def _pdf_page_count(path: Path) -> int | None:
    bundle = _fitz_bundle(path)
    if bundle["ok"] and bundle["page_count"] is not None:
        return int(bundle["page_count"])
    try:
        from pypdf import PdfReader  # type: ignore

        return len(PdfReader(str(path)).pages)
    except Exception:
        return None


def _docx_text_only_words_per_page(word_count: int) -> int:
    """Words-per-page for DOCX with no ``inline_shapes`` (spacing differs by length)."""
    if word_count >= DOCX_TEXT_ONLY_LONG_WORD_THRESHOLD:
        return ESTIMATED_WORDS_PER_PAGE_DOCUMENT_TEXT_ONLY_LONG
    return ESTIMATED_WORDS_PER_PAGE_DOCUMENT_TEXT_ONLY


def _pptx_slide_count(path: Path) -> int | None:
    try:
        from pptx import Presentation  # type: ignore

        return len(Presentation(str(path)).slides)
    except Exception:
        return None


def _docx_page_count_estimate(path: Path, word_count: int) -> int:
    """Words-per-page estimate plus extra pages from tall embedded ``inline_shapes``."""
    try:
        import docx  # type: ignore

        shapes = list(docx.Document(str(path)).inline_shapes)
    except Exception:
        shapes = []

    # Text-only: spacing in Word varies—short docs vs long Notion exports use tiered w/page.
    if not shapes:
        per = _docx_text_only_words_per_page(word_count)
        return _pages_from_word_estimate(word_count, per_page=per)

    word_pages = _pages_from_word_estimate(
        word_count, per_page=ESTIMATED_WORDS_PER_PAGE_DOCUMENT
    )
    sum_h = sum(int(s.height) for s in shapes if s.height)
    if sum_h <= 0:
        return word_pages

    raw_stack_pages = math.ceil(sum_h / DOCX_IMAGE_STACK_PAGE_EMU)
    # One modest inline graphic is often already "paid for" by surrounding text pages.
    if len(shapes) <= 1:
        image_extra = max(0, raw_stack_pages - 1)
    else:
        image_extra = raw_stack_pages
    return word_pages + image_extra


def _xlsx_sheet_count(path: Path) -> int | None:
    try:
        import openpyxl  # type: ignore

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        try:
            n = len(wb.sheetnames)
            return n if n > 0 else None
        finally:
            wb.close()
    except Exception:
        return None


def _xls_sheet_count(path: Path) -> int | None:
    try:
        import xlrd  # type: ignore

        wb = xlrd.open_workbook(str(path))
        n = wb.nsheets
        return n if n > 0 else None
    except Exception:
        return None


def page_count_from_file(
    file_path: str | Path,
    *,
    html_via_docx: bool = True,
    word_count: int | None = None,
) -> int:
    """Return a best-effort page (or sheet/slide) count for inventory display.

    - **PDF**: exact page count from the file.
    - **PPTX**: number of slides.
    - **XLSX/XLSM**: number of worksheets (treated as a multi-page workbook).
    - **DOCX**: tiered **text-only** words/page (see :data:`DOCX_TEXT_ONLY_LONG_WORD_THRESHOLD`);
      **120** w/page when ``inline_shapes`` exist, plus image stack pages.
    - **HTML**: word-density only (figures not measured without a full layout).
    - **Other text-like types**: estimated using :data:`ESTIMATED_WORDS_PER_PAGE`.

    Word rules match :func:`count_words_from_file` when ``word_count`` is omitted.
    Counts are approximate without a layout engine — not identical to Word print preview.
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    if word_count is None:
        word_count = count_words_from_file(path, html_via_docx=html_via_docx)

    if ext == ".pdf":
        n = _pdf_page_count(path)
        if n is not None and n > 0:
            return n
        return _pages_from_word_estimate(word_count)

    if ext == ".pptx":
        n = _pptx_slide_count(path)
        if n is not None and n > 0:
            return n
        return _pages_from_word_estimate(word_count)

    if ext == ".docx":
        return _docx_page_count_estimate(path, word_count)

    if ext in {".xlsx", ".xlsm"}:
        n = _xlsx_sheet_count(path)
        if n is not None and n > 0:
            return max(1, n)
        return _pages_from_word_estimate(word_count)

    if ext == ".xls":
        n = _xls_sheet_count(path)
        if n is not None and n > 0:
            return max(1, n)
        return _pages_from_word_estimate(word_count)

    if ext in {".html", ".htm"}:
        return _pages_from_word_estimate(
            word_count, per_page=ESTIMATED_WORDS_PER_PAGE_DOCUMENT
        )

    if ext in {".txt", ".md", ".rst", ".log", ".csv", ".tsv", ".json"}:
        return _pages_from_word_estimate(word_count)

    return _pages_from_word_estimate(word_count)


def extract_snippet(
    file_path: str,
    *,
    max_chars: int = MAX_CHARS_DEFAULT,
    max_bytes: int = MAX_BYTES_READ_DEFAULT,
) -> str | None:
    """Return a short snippet (or None if not extractable)."""
    path = Path(file_path)
    ext = path.suffix.lower()

    try:
        if ext == ".eml":
            subj = extract_email_subject(path)
            return _truncate(subj, max_chars) if subj else None

        if ext in {".txt", ".md", ".rst", ".log"}:
            return _truncate(_read_text_file(path, max_bytes=max_bytes), max_chars)

        if ext in {".html", ".htm"}:
            raw = _read_text_file(path, max_bytes=max_bytes)
            # HTML exports (e.g., Google Docs) often contain ligature glyphs (ﬁ, ﬂ, …)
            # that can look like "missing letters" downstream unless normalized.
            cleaned = _strip_html(raw)
            cleaned = _normalize_extracted_text(cleaned)
            return _truncate(cleaned, max_chars)

        if ext == ".json":
            raw = _read_text_file(path, max_bytes=max_bytes)
            try:
                obj = json.loads(raw)
                pretty = json.dumps(obj, ensure_ascii=False, indent=2)
                return _truncate(pretty, max_chars)
            except Exception:
                return _truncate(raw, max_chars)

        if ext in {".csv", ".tsv"}:
            raw = _read_text_file(path, max_bytes=max_bytes)
            lines = raw.splitlines()
            head = "\n".join(lines[:80])
            return _truncate(head, max_chars)

        if ext in {".xlsx", ".xlsm"}:
            try:
                import openpyxl  # type: ignore
            except Exception:
                return None
            wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
            try:
                lines: list[str] = []
                for sheetname in wb.sheetnames[:5]:
                    ws = wb[sheetname]
                    sheet_lines: list[str] = []
                    for r_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
                        if r_idx > 20:
                            break
                        cells = [("" if v is None else str(v)) for v in row[:20]]
                        if any(c.strip() for c in cells):
                            sheet_lines.append("\t".join(cells))
                    if sheet_lines:
                        lines.append(f"[Sheet: {sheetname}]")
                        lines.extend(sheet_lines)
                        lines.append("")
                if len(wb.sheetnames) > 5:
                    lines.append(f"[... {len(wb.sheetnames) - 5} more sheet(s)]")
            finally:
                wb.close()
            txt = "\n".join(lines).strip()
            return _truncate(txt, max_chars) if txt else None

        if ext == ".xls":
            try:
                import xlrd  # type: ignore
            except Exception:
                return None
            try:
                wb = xlrd.open_workbook(str(path))
            except Exception:
                return None
            lines = []
            for sheet_idx in range(min(5, wb.nsheets)):
                ws = wb.sheet_by_index(sheet_idx)
                sheet_lines: list[str] = []
                for row_idx in range(min(20, ws.nrows)):
                    cells = []
                    for col_idx in range(min(20, ws.ncols)):
                        ctype = ws.cell_type(row_idx, col_idx)
                        val = ws.cell_value(row_idx, col_idx)
                        if ctype == xlrd.XL_CELL_EMPTY or ctype == xlrd.XL_CELL_ERROR:
                            cells.append("")
                        elif ctype == xlrd.XL_CELL_DATE:
                            try:
                                import datetime
                                dt = xlrd.xldate_as_datetime(val, wb.datemode)
                                cells.append(dt.strftime("%Y-%m-%d"))
                            except Exception:
                                cells.append(str(val))
                        elif ctype == xlrd.XL_CELL_NUMBER:
                            cells.append(str(int(val)) if isinstance(val, float) and val == int(val) else str(val))
                        elif ctype == xlrd.XL_CELL_BOOLEAN:
                            cells.append("True" if val else "False")
                        else:
                            cells.append(str(val).strip())
                    if any(c.strip() for c in cells):
                        sheet_lines.append("\t".join(cells))
                if sheet_lines:
                    lines.append(f"[Sheet: {ws.name}]")
                    lines.extend(sheet_lines)
                    lines.append("")
            if wb.nsheets > 5:
                lines.append(f"[... {wb.nsheets - 5} more sheet(s)]")
            txt = "\n".join(lines).strip()
            return _truncate(txt, max_chars) if txt else None

        if ext == ".pdf":
            raw = _extract_pdf_pages(path, max_pages=5)
            return _truncate(raw, max_chars) if raw else None

        if ext == ".docx":
            try:
                import docx  # type: ignore
            except Exception:
                return None
            doc = docx.Document(str(path))
            paras = [_normalize_extracted_text(p.text) for p in doc.paragraphs if p.text.strip()]
            return _truncate("\n".join(paras[:200]), max_chars) if paras else None

        if ext == ".pptx":
            # Optional dependency; keep graceful fallback.
            try:
                from pptx import Presentation  # type: ignore
            except Exception:
                return None
            prs = Presentation(str(path))
            parts: list[str] = []
            # python-pptx Slides collection doesn't support slicing reliably across versions.
            max_slides = min(20, len(prs.slides))
            for i in range(max_slides):
                slide = prs.slides[i]
                for shape in slide.shapes:
                    text = getattr(shape, "text", None)
                    if isinstance(text, str) and text.strip():
                        parts.append(_normalize_extracted_text(text.strip()))
                if sum(len(p) for p in parts) >= max_chars:
                    break
            return _truncate("\n\n".join(parts), max_chars) if parts else None

        if ext in _IMAGE_DOT_EXT:
            t = _extract_image_plaintext(path)
            return _truncate(t, max_chars) if t else None

        # For unknown/binary formats: no snippet.
        return None
    except Exception:
        return None

